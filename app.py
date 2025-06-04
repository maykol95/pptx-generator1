import streamlit as st
import pandas as pd
import json
import os
import tempfile
import shutil
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from datetime import datetime
from PIL import Image
from io import BytesIO
import time
from concurrent.futures import ThreadPoolExecutor

st.set_page_config(layout="wide")
st.title("Generador de PowerPoint desde Excel")

# Inicializar estado de sesión para almacenar descargas
if "archivos_generados" not in st.session_state:
    st.session_state.archivos_generados = []
if "start_time" not in st.session_state:
    st.session_state.start_time = None

def calcular_dimensiones(img_path, fotos_por_slide, encabezados_count):
    try:
        with Image.open(img_path) as img:
            width_px, height_px = img.size
            width_in, height_in = width_px / 96, height_px / 96
            max_width = 11 / fotos_por_slide - 0.5
            encabezado_space = 0.3 + 0.25 * min(encabezados_count, 6)
            max_height = 5.3 - encabezado_space
            ratio = min(max_width / width_in, max_height / height_in)
            return Inches(width_in * ratio), Inches(height_in * ratio)
    except Exception as e:
        print(f"Error al procesar {img_path}: {e}")
        return None, None

excel_file = st.sidebar.file_uploader("Cargar archivo Excel", type=[".xlsx", ".xls"])
st.sidebar.subheader("Imagen de fondo para cada diapositiva")
fondo_file = st.sidebar.file_uploader("Cargar imagen (.jpg o .png)", type=["jpg", "jpeg", "png"])
fondo_bytes = fondo_file.read() if fondo_file else None

st.sidebar.subheader("Estilo de los encabezados")
color_fuente = st.sidebar.color_picker("Color de fuente", "#000000")
tipo_fuente = st.sidebar.selectbox("Tipo de letra", ["Arial", "Calibri", "Times New Roman", "Verdana", "Helvetica"])

if excel_file:
    df = pd.read_excel(excel_file)
    st.subheader("Configuración de la presentación")
    encabezados_seleccionados = st.multiselect("Selecciona columnas como encabezados de imagen (máx 8)", df.columns.tolist(), max_selections=8)
    fotos_por_slide = st.number_input("¿Cuántas fotos por diapositiva? (máx 4)", min_value=1, max_value=4, value=1, step=1)
    orden_columna = st.selectbox("Ordenar por columna (opcional)", [""] + df.columns.tolist())
    subdiv_col = st.selectbox("Subdividir por columna (opcional, genera un archivo por valor distinto)", [""] + df.columns.tolist())

    st.subheader("¿Deseas aplicar filtros?")
    aplicar_filtros = st.checkbox("Sí, mostrar opciones de filtrado")
    if aplicar_filtros:
        columnas = df.columns.tolist()
        filtros = {}

        if 'fecha' in df.columns:
            try:
                df['fecha'] = pd.to_datetime(df['fecha'].astype(str), errors='coerce')
                df = df.dropna(subset=['fecha'])
                fecha_min = df['fecha'].min().date()
                fecha_max = df['fecha'].max().date()
                fecha_inicio, fecha_fin = st.date_input("Selecciona el rango de fechas", value=(fecha_min, fecha_max))
                df = df[(df['fecha'] >= pd.to_datetime(fecha_inicio)) & (df['fecha'] <= pd.to_datetime(fecha_fin))]
            except: pass

        for col in columnas:
            dtype = df[col].dtype
            opciones = df[col].dropna().unique().tolist()
            if dtype == object or dtype.name == 'category':
                seleccion = st.multiselect(f"Filtrar por {col}", sorted(map(str, opciones)))
                if seleccion:
                    filtros[col] = seleccion
            elif pd.api.types.is_numeric_dtype(dtype):
                valores_input = st.text_input(f"Valores para {col} (separados por comas)")
                if valores_input:
                    try:
                        valores = [float(v.strip()) for v in valores_input.split(',')]
                        filtros[col] = valores
                    except: pass
            elif pd.api.types.is_bool_dtype(dtype):
                seleccion = st.multiselect(f"Filtrar por {col}", [True, False])
                if seleccion:
                    filtros[col] = seleccion

        for col, valores in filtros.items():
            df = df[df[col].isin(valores)]

        st.write(f"Filas después del filtrado: {len(df)}")

    if st.button("Generar PowerPoint"):
        st.session_state.start_time = datetime.now()
        progress_bar = st.progress(0)
        status_text = st.empty()

        try:
            if orden_columna:
                df = df.sort_values(by=orden_columna)

            temp_dir = tempfile.mkdtemp()
            imagen_col = None
            for col in df.columns:
                if df[col].astype(str).str.contains("https?://.*\\.(jpg|jpeg|png)", case=False, na=False).any():
                    imagen_col = col
                    break

            if not imagen_col:
                st.warning("No se encontró ninguna columna con URLs de imagen válidas.")
            else:
                def descargar_imagen(row, i):
                    url = row[imagen_col]
                    try:
                        response = requests.get(url, stream=True, timeout=5)
                        if response.status_code == 200:
                            img_path = os.path.join(temp_dir, f"img_{i}.jpg")
                            with open(img_path, 'wb') as f:
                                f.write(response.content)
                            if os.path.getsize(img_path) > 0:
                                row['img_path'] = img_path
                                return row.to_dict()
                    except:
                        return None
                    return None

                with ThreadPoolExecutor(max_workers=16) as executor:
                    futures = [executor.submit(descargar_imagen, row, i) for i, (_, row) in enumerate(df.iterrows())]
                    resultados = [f.result() for f in futures if f.result()]

                df_filtrado = pd.DataFrame(resultados)

                def generar_presentacion(df_slice, nombre):
                    prs = Presentation()
                    prs.slide_width = Inches(13.33)
                    prs.slide_height = Inches(7.5)
                    layout = prs.slide_layouts[6]

                    if fondo_bytes:
                        fondo_img = Image.open(BytesIO(fondo_bytes))
                        fondo_path = os.path.join(temp_dir, "fondo_temp.jpg")
                        fondo_img.save(fondo_path)

                    for i in range(0, len(df_slice), fotos_por_slide):
                        slide = prs.slides.add_slide(layout)
                        if fondo_bytes:
                            slide.shapes.add_picture(fondo_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

                        subset = df_slice.iloc[i:i+fotos_por_slide]
                        spacing_x = Inches(12 / fotos_por_slide)
                        for j, (_, row) in enumerate(subset.iterrows()):
                            img_path = row['img_path']
                            img_width, img_height = calcular_dimensiones(img_path, fotos_por_slide, len(encabezados_seleccionados))
                            img_left = spacing_x * j + Inches(1.15)
                            encabezado_height = Inches(0.25 * min(len(encabezados_seleccionados), 6))
                            img_top = encabezado_height + Inches(0.4)

                            if img_width and img_height:
                                slide.shapes.add_picture(img_path, img_left, img_top, width=img_width, height=img_height)
                                encabezado = " | ".join(f"{col}: {row.get(col, '')}" for col in encabezados_seleccionados if col in row)
                                if encabezado:
                                    text_box = slide.shapes.add_textbox(img_left, Inches(0.2), img_width, encabezado_height)
                                    text_frame = text_box.text_frame
                                    text_frame.clear()
                                    for line in encabezado.split(" | "):
                                        p = text_frame.add_paragraph()
                                        p.text = line
                                        p.font.size = Pt(8)
                                        p.font.bold = True
                                        p.font.name = tipo_fuente
                                        r, g, b = [int(color_fuente.lstrip("#")[i:i+2], 16) for i in (0, 2, 4)]
                                        p.font.color.rgb = RGBColor(r, g, b)

                    path = os.path.join(temp_dir, f"{nombre}.pptx")
                    prs.save(path)
                    return path

                archivos_generados = []
                if subdiv_col and subdiv_col in df_filtrado.columns:
                    valores_unicos = df_filtrado[subdiv_col].dropna().unique()
                    total = len(valores_unicos)
                    for idx, val in enumerate(valores_unicos):
                        status_text.text(f"Generando presentación para '{val}' ({idx + 1}/{total})")
                        df_part = df_filtrado[df_filtrado[subdiv_col] == val]
                        pptx_path = generar_presentacion(df_part, str(val))
                        archivos_generados.append((str(val), pptx_path))
                        progress_bar.progress((idx + 1) / total)
                else:
                    status_text.text("Generando presentación completa")
                    pptx_path = generar_presentacion(df_filtrado, "presentacion")
                    archivos_generados.append(("presentacion", pptx_path))
                    progress_bar.progress(1.0)

                st.session_state.archivos_generados = archivos_generados

        except Exception as e:
            st.error(f"Error al generar presentación: {e}")

# Mostrar descargas y tiempo
if st.session_state.archivos_generados:
    st.subheader("Descargas disponibles")
    for nombre, path in st.session_state.archivos_generados:
        with open(path, "rb") as f:
            st.download_button(
                label=f"⬇ Descargar {nombre}.pptx",
                data=f,
                file_name=f"{nombre}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentation.presentation"
            )
    if st.session_state.start_time:
        elapsed = datetime.now() - st.session_state.start_time
        st.success(f"Presentación generada en {elapsed.total_seconds():.2f} segundos.")
